# -*- coding: utf-8 -*-
"""
HSR 作品集演示 PPT 生成器
面向面试/作品集的 17 页演示，全中文，画面与图示为主，不放代码/IDE。
深色 HSR 风格：近黑背景 #14171d，青色强调 #6fd7ff，金色点缀 #ffd479。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 常量 ─────────────────────────────────────────────────────────────
BG      = RGBColor(0x14, 0x17, 0x1d)   # 近黑
PANEL   = RGBColor(0x1d, 0x22, 0x2b)   # 卡片底色
PANEL2  = RGBColor(0x25, 0x2b, 0x36)   # 卡片描边/次级
CYAN    = RGBColor(0x6f, 0xd7, 0xff)   # 强调青
GOLD    = RGBColor(0xff, 0xd4, 0x79)   # 点缀金
WHITE   = RGBColor(0xf2, 0xf3, 0xf5)
GRAY    = RGBColor(0x9a, 0xa3, 0xad)
DIM     = RGBColor(0x6c, 0x75, 0x80)
GREEN   = RGBColor(0x7f, 0xd9, 0x8a)   # 完成
BLUE    = RGBColor(0x7a, 0xb8, 0xff)   # 进行中

ASSET_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "ppt_assets")
OUT_PATH  = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "HSR_Demo_Presentation_v2.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── 工具函数 ─────────────────────────────────────────────────────────
def _set_text(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tf.clear()
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Microsoft YaHei"
    return tf

def add_text(slide, x, y, w, h, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    _set_text(tf, text, size, color, bold, align)
    return tb

def add_rect(slide, x, y, w, h, fill=PANEL, line=None, line_w=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    if radius:
        # 调整圆角半径
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    shp.shadow.inherit = False
    return shp

def add_card(slide, x, y, w, h, title, body, title_color=CYAN, fill=PANEL, title_size=16, body_size=12):
    add_rect(slide, x, y, w, h, fill=fill, line=PANEL2, radius=True)
    add_text(slide, x + Inches(0.25), y + Inches(0.2), w - Inches(0.5), Inches(0.5),
             title, size=title_size, color=title_color, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(0.75), w - Inches(0.5), h - Inches(0.95),
             body, size=body_size, color=GRAY)

def add_pill(slide, x, y, w, h, text, fill, color=WHITE, size=12, bold=True):
    shp = add_rect(slide, x, y, w, h, fill=fill, radius=True)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_top = tf.margin_bottom = 0
    _set_text(tf, text, size, color, bold, PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shp

def add_arrow(slide, x, y, w=Inches(0.9), h=Inches(0.5), color=CYAN, flip=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    if flip:
        shp.rotation = 180
    return shp

def add_header(slide, kicker, title, page_no):
    """顶部页眉：左上小标签 + 标题，右下页码"""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill=CYAN)
    add_text(slide, Inches(0.55), Inches(0.3), Inches(8), Inches(0.3),
             kicker, size=11, color=CYAN, bold=True)
    add_text(slide, Inches(0.55), Inches(0.55), Inches(9), Inches(0.7),
             title, size=28, color=WHITE, bold=True)
    add_text(slide, SLIDE_W - Inches(1.2), Inches(0.4), Inches(0.7), Inches(0.4),
             f"{page_no:02d} / 17", size=11, color=DIM, align=PP_ALIGN.RIGHT)

def add_image_fit(slide, img_path, x, y, w, h, border=True):
    """按比例放入图片，保持纵横比居中"""
    if not os.path.exists(img_path):
        add_rect(slide, x, y, w, h, fill=PANEL2, line=PANEL2)
        add_text(slide, x, y, w, h, "截图待补充\n(录制讲解时截取)", size=12, color=DIM,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return
    from PIL import Image
    iw, ih = Image.open(img_path).size
    ar = iw / ih
    box_ar = (w / h)
    if ar > box_ar:
        nw = w; nh = int(w / ar)
    else:
        nh = h; nw = int(h * ar)
    ox = x + int((w - nw) / 2); oy = y + int((h - nh) / 2)
    pic = slide.shapes.add_picture(img_path, ox, oy, nw, nh)
    if border:
        add_rect(slide, ox, oy, nw, nh, fill=None, line=CYAN, line_w=Pt(1.5))
    return pic

def add_footer_note(slide, note):
    add_text(slide, Inches(0.55), SLIDE_H - Inches(0.45), SLIDE_W - Inches(1.1), Inches(0.3),
             note, size=10, color=DIM, align=PP_ALIGN.CENTER)

# ── 幻灯片构建 ───────────────────────────────────────────────────────
def new_slide():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=BG)
    return s

def flow_bar(slide, y, labels, current, colors):
    """横向流程条：labels 数组 + 当前步高亮。colors: 每步颜色数组"""
    n = len(labels)
    gap = Inches(0.25)
    total_gap = gap * (n - 1)
    seg_w = (Inches(12.2) - total_gap) / n
    x = Inches(0.55)
    for i, lab in enumerate(labels):
        c = colors[i]
        seg = add_rect(slide, x, y, seg_w, Inches(0.62), fill=c if i <= current else PANEL2, radius=True)
        tf = seg.text_frame; tf.word_wrap = False
        _set_text(tf, lab, 13 if i == current else 12, WHITE, bold=(i == current), align=PP_ALIGN.CENTER)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        x = x + seg_w + gap

def stat_tile(slide, x, y, w, h, number, label, color=CYAN, number_size=30):
    add_rect(slide, x, y, w, h, fill=PANEL, line=PANEL2, radius=True)
    add_text(slide, x, y + Inches(0.12), w, Inches(0.65), number, size=number_size, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + h - Inches(0.62), w, Inches(0.5), label, size=11, color=GRAY, align=PP_ALIGN.CENTER)

def check_list(slide, x, y, w, items, size=13, gap=0.32, color=GRAY, mark_color=GREEN):
    yy = y
    for it in items:
        add_text(slide, x, yy, Inches(0.4), Inches(0.35), "✓", size=size, color=mark_color, bold=True)
        add_text(slide, x + Inches(0.4), yy, w - Inches(0.4), Inches(0.35), it, size=size, color=color)
        yy += Inches(gap)

# ═══════════════════════════════════════════════════════════════════
# P1 封面
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
cover_img = os.path.join(ASSET_DIR, "exploration_scrapopolis.png")
if os.path.exists(cover_img):
    add_image_fit(s, cover_img, Inches(0), Inches(0), SLIDE_W, SLIDE_H, border=False)
# 半透明遮罩
ov = add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=RGBColor(0x0c, 0x0f, 0x14))
ov.fill.fore_color.rgb = RGBColor(0x0c, 0x0f, 0x14)
try: ov.fill.transparency = 40
except Exception: pass
add_pill(s, Inches(0.9), Inches(1.15), Inches(2.6), Inches(0.45), "作品集 · 演示", CYAN, WHITE, 13)
add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.6), "HSR", size=88, color=WHITE, bold=True)
add_text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.8),
         "从零搭建的探索 × 回合制 JRPG", size=30, color=CYAN, bold=True)
add_text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.5),
         "Unreal Engine 5.6 · C++20 · Gameplay Ability System · UMG", size=16, color=GRAY)
add_text(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.5),
         "原创作 · 从空工程搭建完整可玩闭环 · 探索 → 回合战斗 → 养成 → 存档", size=14, color=GRAY)

# ═══════════════════════════════════════════════════════════════════
# P2 项目一句话 + 目标
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "PROJECT OVERVIEW", "项目一句话", 2)
add_rect(s, Inches(0.55), Inches(1.75), Inches(12.2), Inches(1.35), fill=PANEL, line=PANEL2, radius=True)
add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.9),
         "用 UE5.6 + GAS 从空白 C++ 工程，一步步搭建出「探索 → 遭遇 → 回合战斗 → 结算 → 养成 → 存档」的原创单机 JRPG 完整闭环。",
         size=19, color=WHITE)
# 三个目标卡片
goals = [
    ("完整闭环", "能从头玩到尾\n探索 / 战斗 / 成长 / 存档\n数据全部可持久化"),
    ("工程严谨", "TDD + 自动化测试\n纯值 DTO 跨系统边界\n权威分离、失败可恢复"),
    ("原创内容", "正式地图 / 角色 / 技能 / 遗器\n全部重新作者化\n不复制任何商业作品"),
]
gx = Inches(0.55); gw = Inches(3.9); gap = Inches(0.25)
for i, (t, d) in enumerate(goals):
    cx = gx + (gw + gap) * i
    add_card(s, cx, Inches(3.6), gw, Inches(2.6), t, d, title_color=[CYAN, GOLD, GREEN][i])
add_footer_note(s, "说明：工程严谨度是本项目区别于普通“作业级 Demo”的核心差异点。")

# ═══════════════════════════════════════════════════════════════════
# P3 制作路线全景
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "ROADMAP", "制作路线 · 18 个阶段从零推进", 3)
# 阶段条（已完成的亮，当前进行中闪烁）
phases = [
    ("P0\n工程基线", True), ("P1\n探索角色", True), ("P2\nGAS 基础", True),
    ("P3\n交互系统", True), ("P4\n敌人/遭遇", True), ("P5\n回合战斗", True),
    ("P6\n技能系统", True), ("P7\n伤害公式", True), ("P8\n弱点/击破", True),
    ("P9\n状态效果", True), ("P10\n战斗UI", True), ("P11\n角色成长", True),
    ("P12\n装备遗器", True), ("P13\n背包奖励", True), ("P14\n任务对话", True),
    ("P15\n地图传送", True), ("P16\n存档系统", True), ("P17\nUI整合", True),
    ("P18\n正式内容", False),
]
n = len(phases)
seg_w = Inches(0.58); gapx = Inches(0.06)
x = Inches(0.45)
for i, (lab, done) in enumerate(phases):
    c = GREEN if done else CYAN
    seg = add_rect(s, x, Inches(1.5), seg_w, Inches(1.05), fill=c if done else PANEL2, radius=True)
    tf = seg.text_frame; tf.word_wrap = True
    _set_text(tf, lab, 8.5, WHITE, bold=done, align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    x += seg_w + gapx
add_text(s, Inches(0.45), Inches(2.75), Inches(12.2), Inches(0.4),
         "P0–P17 系统层全部完成 ✓    当前：正式 Demo 内容整合（地图 / 角色 / 技能 / 遗器）", size=14, color=GRAY)
add_rect(s, Inches(0.55), Inches(3.5), Inches(12.2), Inches(3.3), fill=PANEL, line=PANEL2, radius=True)
add_text(s, Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.4), "内容主线（每个阶段一个可玩闭环）", size=16, color=CYAN, bold=True)
flow_bar(s, Inches(4.4), ["探索", "遭遇", "回合战斗", "养成", "存档"], 2,
         [GREEN, GREEN, CYAN, GREEN, GREEN])
add_text(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.9),
         "· 探索移动、交互、敌人感知/追击、三张正式地图传送\n· 独立战斗地图内速度排序回合制、技能/击破/状态、胜负结算\n· 角色成长、遗器/装备、背包/奖励/掉落\n· 全量权威数据存档、版本迁移、冷启动恢复", size=13, color=GRAY)
add_footer_note(s, "阶段推进采用 TDD：先写失败测试 → 实现 → 构建 + 自动化回归 → 用户 PIE 验收。")

# ═══════════════════════════════════════════════════════════════════
# P4 玩法闭环总览
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "GAMEPLAY LOOP", "玩法闭环总览", 4)
loop_nodes = ["探索移动\n三张地图 · 传送", "遭遇敌人\n感知 → 追击 → 触发", "回合战斗\n速度条 · 技能 · 击破", "结算奖励\n经验 · 掉落 · 遗器", "养成\n角色 · 遗器强化", "存档 / 读档\n全量持久化"]
n = len(loop_nodes)
seg_w = Inches(1.95); gapx = Inches(0.12)
x = Inches(0.5); y = Inches(1.9)
for i, lab in enumerate(loop_nodes):
    add_card(s, x, y, seg_w, Inches(1.5), f"0{i+1}", lab, title_color=CYAN, title_size=18)
    if i < n - 1:
        add_arrow(s, x + seg_w + Inches(0.01), y + Inches(0.5), Inches(0.1), Inches(0.5))
    x += seg_w + gapx + Inches(0.1)
# 探索图
add_text(s, Inches(0.55), Inches(3.7), Inches(6), Inches(0.4), "探索画面（Scrapopolis 正式场景）", size=15, color=CYAN, bold=True)
add_image_fit(s, os.path.join(ASSET_DIR, "exploration_scrapopolis.png"), Inches(0.55), Inches(4.15), Inches(6.2), Inches(2.6))
# 闭环文字
add_card(s, Inches(7.2), Inches(3.7), Inches(5.6), Inches(3.05), "一次完整循环",
         "从探索地图出发 → 遭遇敌人进入独立战斗地图 → 回合制分胜负 → 胜利结算经验与掉落 → 在背包/遗器面板强化成长 → 保存后任意时刻读档继续。\n\n系统之间只传纯值数据快照，没有跨地图的 Actor 引用，因此存档、回放和自动化测试都安全可靠。",
         title_color=GOLD)
add_footer_note(s, "这张闭环图是全篇主线：之后每页都是其中一环的实现细节。")

# ═══════════════════════════════════════════════════════════════════
# P5 探索层
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "EXPLORATION", "探索层 · 第一人称到第三人称的完整移动", 5)
# 左：要点
add_card(s, Inches(0.55), Inches(1.5), Inches(5.6), Inches(5.3), "实现要点",
         "· 第三人称移动 / 镜头 / 跳跃（Enhanced Input 数据驱动）\n\n"
         "· 交互系统：近距对象检测、F 键交互、事件驱动观察层\n\n"
         "· 敌人 AI：巡逻 → 感知 → 追击 → 丢失返回（Behavior Tree / Blackboard）\n\n"
         "· 三张正式探索地图互通：\n  观景车厢 / 新艾丽都六分街 / 黑塔空间站\n\n"
         "· 地图面板数据驱动：按钮显示目标地图中文名\n  巡逻敌人进战 + 多个内容不同的宝箱")
# 右：两张图
add_text(s, Inches(6.4), Inches(1.5), Inches(6.4), Inches(0.35), "正式场景探索", size=14, color=CYAN, bold=True)
add_image_fit(s, os.path.join(ASSET_DIR, "exploration_scrapopolis.png"), Inches(6.4), Inches(1.9), Inches(6.35), Inches(2.5))
add_text(s, Inches(6.4), Inches(4.55), Inches(6.4), Inches(0.35), "早期灰盒验证（探索地图 A）", size=12, color=DIM)
add_image_fit(s, os.path.join(ASSET_DIR, "exploration_graybox_b.png"), Inches(6.4), Inches(4.9), Inches(6.35), Inches(1.9))
add_footer_note(s, "录屏建议：展示移动 + 走进 NPC 出现“交互”提示 + 敌人发现玩家开始追击的一段。")

# ═══════════════════════════════════════════════════════════════════
# P6 回合制战斗（核心页）
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "TURN-BASED BATTLE", "回合制战斗 · 速度驱动的指令系统", 6)
# 上半：流程
flow_bar(s, Inches(1.5), ["行动条排序", "玩家指令", "目标选择", "伤害结算", "回合推进"], 1,
         [GREEN, CYAN, GREEN, GREEN, GREEN])
# 左卡片：系统
add_card(s, Inches(0.55), Inches(2.5), Inches(4.0), Inches(4.3), "战斗系统",
         "· 速度行动条：无 Tick 事件驱动\n"
         "· 普攻 / 战技 / 终结技 / 治疗，SP 与能量管理\n"
         "· 弱点 → 削韧 → 击破 → 行动延后\n"
         "· 状态效果：Buff / DoT / 免疫 / 驱散\n"
         "· 原创伤害公式 + 暴击 / 防御 / 属性克制",
         title_color=CYAN)
# 中卡片：指令面板（占位截图）
add_card(s, Inches(4.8), Inches(2.5), Inches(4.0), Inches(4.3), "指令面板",
         "（此处放战斗指令面板截图）\n\n角色头像 / SP / 能量\n技能按钮（技能名 + 消耗）\n目标选择高亮", title_color=GOLD)
# 右卡片：伤害链路
add_card(s, Inches(9.05), Inches(2.5), Inches(3.75), Inches(4.3), "一次伤害怎么算",
         "选技能\n→ 消耗 SP / 能量\n→ 判定命中与暴击\n→ 计算防御减伤\n→ 输出伤害 + 削韧\n→ 状态效果触发\n→ 结算事件发布",
         title_color=GREEN)
add_footer_note(s, "录屏建议：完整展示一次战斗——选技 → 目标 → 结算 → 敌人回合，配讲解。")

# ═══════════════════════════════════════════════════════════════════
# P7 养成系统
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "PROGRESSION", "养成系统 · 角色成长 × 遗器装备", 7)
# 三列卡片
cards = [
    ("角色成长", "等级 / 经验曲线\n属性随等级成长\n技能等级解锁\n（角色详情面板）", CYAN),
    ("遗器系统", "套装（2 件套 / 4 件套）\n强化 / 词条 / 部位\n属性来源可逆追踪\n卸下即还原", GOLD),
    ("背包 / 奖励", "物品分类 / 筛选 / 排序\n掉落与奖励结算\n重复奖励防重\n资源上限保护", GREEN),
]
cx = Inches(0.55); cw = Inches(4.0); cgap = Inches(0.1)
for i, (t, d, c) in enumerate(cards):
    add_card(s, cx + (cw + cgap) * i, Inches(1.5), cw, Inches(2.7), t, d, title_color=c)
add_text(s, Inches(0.55), Inches(4.5), Inches(12.2), Inches(0.35), "遗器装备面板（截图待补充）", size=14, color=CYAN, bold=True)
add_image_fit(s, "", Inches(0.55), Inches(4.9), Inches(6.1), Inches(2.1))
add_text(s, Inches(6.9), Inches(4.5), Inches(6.0), Inches(0.35), "背包 / 遗器强化面板（截图待补充）", size=14, color=CYAN, bold=True)
add_image_fit(s, "", Inches(6.9), Inches(4.9), Inches(5.9), Inches(2.1))
add_footer_note(s, "录屏建议：展示角色详情属性、遗器装备前后属性变化、背包分类筛选。")

# ═══════════════════════════════════════════════════════════════════
# P8 任务与对话
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "QUEST & DIALOGUE", "任务与对话 · 「界域回响」", 8)
# 左：任务系统
add_card(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(5.3), "任务系统（Quest Subsystem）",
         "· 正式任务「界域回响」已作者化：\n"
         "  观景车厢调查 → 六分街调查 → 开启宝箱\n  → 击败巡检机 → 击败 Boss 来古士\n\n"
         "· 目标事件驱动：任务进度随探索 / 战斗事件推进\n\n"
         "· 任务与对话 / 遭遇 / 奖励联动：\n"
         "  对话选项可触发遭遇，胜利结算任务与奖励\n\n"
         "· 状态持久化：任务进度进入存档",
         title_color=CYAN)
# 右：对话系统 + 截图占位
add_card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.5), "对话系统（NPC 分支对话）",
         "· NPC（凯瑟琳）对话，分支选项\n· 对话可推进任务 / 触发遭遇 / 发放奖励\n· Overlay 独立于暂停菜单，交互自然", title_color=GOLD)
add_text(s, Inches(6.8), Inches(4.25), Inches(6.0), Inches(0.35), "对话 Overlay 截图（待补充）", size=13, color=CYAN, bold=True)
add_image_fit(s, "", Inches(6.8), Inches(4.6), Inches(6.0), Inches(2.2))
add_footer_note(s, "诚实说明：任务数据不暴露标题字段，「界域回响」标题出现在任务 UI 仍需后续呈现字段任务。")

# ═══════════════════════════════════════════════════════════════════
# P9 存档系统
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "SAVE SYSTEM", "存档系统 · 全量权威数据持久化", 9)
# 中间示意：系统 → 存档
sys_labels = ["角色 / 队伍", "背包 / 遗器", "任务进度", "地图 / 位置", "奖励账本"]
add_text(s, Inches(0.55), Inches(1.5), Inches(6), Inches(0.35), "存档覆盖的权威数据", size=15, color=CYAN, bold=True)
sy = Inches(2.0)
for i, lab in enumerate(sys_labels):
    add_pill(s, Inches(0.55), sy, Inches(3.2), Inches(0.5), lab, PANEL, WHITE, 13)
    add_arrow(s, Inches(3.9), sy + Inches(0.05), Inches(0.7), Inches(0.4))
    add_pill(s, Inches(4.7), sy, Inches(3.2), Inches(0.5), "存 / 读 快照", [GREEN if i % 2 == 0 else GOLD][0], WHITE, 13)
    sy += Inches(0.66)
# 右：特性
add_card(s, Inches(8.3), Inches(1.5), Inches(4.5), Inches(5.3), "可靠性设计",
         "· 事务写入：先写暂存 → 提交主档 → 校验\n\n"
         "· 版本化 Schema + 保守迁移\n  旧档升级不丢数据\n\n"
         "· 冷启动恢复：任意阶段读档\n  队伍 / 位置 / 任务 / 奖励全部还原\n\n"
         "· 失败可恢复：断电 / 写坏档有回退\n\n"
         "· 全部使用纯值 DTO 序列化\n  不保存 Actor / Widget 引用",
         title_color=GOLD)
add_footer_note(s, "录屏建议：进入游戏 → 做几件事 → 保存 → 退出 → 重新打开 → 状态完整恢复。")

# ═══════════════════════════════════════════════════════════════════
# P10 UI 总整合
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "UI INTEGRATION", "UI 总整合 · 深色 HSR 风格前端", 10)
# 顶部：暂停中枢
add_card(s, Inches(0.55), Inches(1.45), Inches(12.2), Inches(1.15), "暂停中枢（Pause Hub）",
         "Tab / 数字键打开，统一深色风格。所有模块通过同一套 输入策略 + 焦点导航 + 返回逻辑 接入。",
         title_color=CYAN)
# 模块矩阵
modules = [
    ("角色", "角色详情 / 养成"),
    ("队伍", "固定双槽 / 编队"),
    ("背包", "分类 / 详情 / 装备"),
    ("遗器", "套装 / 强化"),
    ("任务", "目标进度"),
    ("地图", "地图 / 传送"),
    ("挑战", "目录 / 战前编队"),
    ("存档", "存 / 读"),
]
mx = Inches(0.55); my = Inches(2.9)
cw = Inches(2.98); ch = Inches(1.35); mgx = Inches(0.1); mgy = Inches(0.15)
for i, (t, d) in enumerate(modules):
    col = i % 4; row = i // 4
    x = mx + (cw + mgx) * col; y = my + (ch + mgy) * row
    add_card(s, x, y, cw, ch, t, d, title_color=[CYAN, GOLD, GREEN, BLUE][i % 4], title_size=14, body_size=11)
add_footer_note(s, "每个面板一张截图（待补充）：角色 / 队伍 / 背包 / 遗器 / 任务 / 地图 / 挑战 / 存档。")

# ═══════════════════════════════════════════════════════════════════
# P11 正式 Demo 内容
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "DEMO CONTENT", "正在制作 · 正式 Demo 内容", 11)
add_text(s, Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.5),
         "已冻结正式内容目录，Asset 校验门禁已 GREEN。全部为原创作者化内容，不复用测试资产。",
         size=15, color=GRAY)
# 三张地图
add_text(s, Inches(0.55), Inches(2.05), Inches(6), Inches(0.35), "三张正式地图", size=15, color=CYAN, bold=True)
maps = [("观景车厢", "探索地图 A"), ("新艾丽都六分街地铁站", "探索地图 B"), ("黑塔空间站支援舱段", "探索地图 C")]
mx = Inches(0.55); mw = Inches(3.9); mgap = Inches(0.25)
for i, (nm, sub) in enumerate(maps):
    add_card(s, mx + (mw + mgap) * i, Inches(2.45), mw, Inches(1.0), nm, sub, title_color=[CYAN, GOLD, GREEN][i], title_size=15)
# 四角色
add_text(s, Inches(0.55), Inches(3.65), Inches(6), Inches(0.35), "四名原创角色 · 12 个技能", size=15, color=CYAN, bold=True)
chars = [("长夜月", "主C · 电弧"), ("火花", "DoT · 疾风"), ("蕾米埃尔", "击破 · 电弧"), ("维林娜", "治疗 · 潮汐")]
cx = Inches(0.55); cw = Inches(2.95); cgap = Inches(0.1)
for i, (nm, role) in enumerate(chars):
    add_card(s, cx + (cw + cgap) * i, Inches(4.1), cw, Inches(0.95), nm, role, title_color=[CYAN, GOLD, GREEN, BLUE][i], title_size=14, body_size=11)
# 底部：遗器 + 任务 + 敌人
add_text(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(0.35), "遗器与内容", size=15, color=CYAN, bold=True)
add_card(s, Inches(0.55), Inches(5.75), Inches(3.95), Inches(1.15), "2 套遗器",
         "闪耀功勋的魔法少女（4 件套）\n天国@直播间（2 件套）", title_color=GOLD, title_size=13, body_size=11)
add_card(s, Inches(4.75), Inches(5.75), Inches(3.95), Inches(1.15), "任务「界域回响」",
         "5 个目标事件驱动\n奖励防重", title_color=CYAN, title_size=13, body_size=11)
add_card(s, Inches(8.95), Inches(5.75), Inches(3.8), Inches(1.15), "敌人与 Boss",
         "支援舱段巡检机\nBoss：来古士", title_color=GREEN, title_size=13, body_size=11)
add_footer_note(s, "角色 / 场景美术为占位，正式角色立绘与场景美术由你后续作者化。")

# ═══════════════════════════════════════════════════════════════════
# P12 架构亮点（图不用代码）
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "ARCHITECTURE", "架构亮点 · 用图不用代码", 12)
# 分层示意
layers = [
    ("表现层 UI", "UMG 只订阅纯值快照", CYAN),
    ("表现事件", "Damage / Toughness / Break / Heal 事件", BLUE),
    ("战斗权威", "Coordinator 单一权威\nTurnManager / Targeting / Status", GOLD),
    ("数据层", "DataAsset 定义 / 纯值 DTO", GREEN),
]
ly = Inches(1.5); lh = Inches(1.15); lgap = Inches(0.12)
for i, (t, d, c) in enumerate(layers):
    y = ly + (lh + lgap) * i
    add_rect(s, Inches(0.55), y, Inches(12.2), lh, fill=PANEL, line=c, radius=True)
    add_text(s, Inches(0.9), y + Inches(0.15), Inches(3.2), Inches(0.5), t, size=17, color=c, bold=True)
    add_text(s, Inches(4.3), y + Inches(0.2), Inches(8.2), Inches(0.8), d, size=13, color=GRAY)
    if i < len(layers) - 1:
        add_text(s, Inches(6.0), y + lh - Inches(0.02), Inches(1.5), Inches(0.18), "▼ 纯值传递", size=9, color=DIM, align=PP_ALIGN.CENTER)
# 右侧要点（叠在下方）
add_card(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.9), "", "核心原则：系统之间只传纯值结构体（DTO），不传 Actor / Widget / GE 引用 → 回放安全、存档安全、测试安全。",
         title_color=CYAN, title_size=11, body_size=13)
add_footer_note(s, "这是面试重点：能讲清“为什么用纯值 DTO”比讲“用了什么库”重要得多。")

# ═══════════════════════════════════════════════════════════════════
# P13 项目数据流动图
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "DATA FLOW", "数据怎么流动 · 一次完整游戏循环", 13)
flow_nodes = [
    ("探索地图", "移动 / 交互 / 敌人感知\n追击触发遭遇", CYAN),
    ("遭遇", "FHSREncounterRequest\n纯值请求 · 世界切换", BLUE),
    ("回合战斗", "速度排序 → 指令 → 结算\nCoordinator 单一权威", GOLD),
    ("结算 / 养成", "经验 · 掉落 · 遗器\n写回运行时数据", GREEN),
    ("存档", "FHSRSaveData 快照\n事务写入 · 版本迁移", GRAY),
]
fx = Inches(0.55); fw = Inches(2.32); fgap = Inches(0.18)
for i, (t, d, c) in enumerate(flow_nodes):
    x = fx + (fw + fgap) * i
    add_card(s, x, Inches(1.5), fw, Inches(1.5), t, d, title_color=c, title_size=16, body_size=11)
    if i < 4:
        add_arrow(s, x + fw + Inches(0.01), Inches(2.0), Inches(0.16), Inches(0.45))
add_rect(s, Inches(0.55), Inches(3.25), Inches(12.23), Inches(0.62), fill=PANEL, line=CYAN, radius=True)
add_text(s, Inches(0.85), Inches(3.35), Inches(11.6), Inches(0.45),
         "跨地图、跨系统的每一次握手都只传纯值结构体：DTO 里没有 Actor / Widget / GameplayEffect 引用",
         size=14, color=CYAN, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.55), Inches(4.1), Inches(12), Inches(0.35), "从探索到存档，数据一路经过的边界", size=15, color=WHITE, bold=True)
detail = [
    ("① 探索 → 遭遇", "AHSREnemyCharacter 重叠 → 构造 FHSREncounterRequest（纯值请求，含队伍 / 遭遇 ID）", CYAN),
    ("② 遭遇 → 战斗", "UHSRBattleTransitionSubsystem 携带请求跨世界，AHSRBattleGameMode 读取", BLUE),
    ("③ 战斗内部", "FHSRBattleParticipant / FHSRDamageResult / FHSRStatusInstance —— 全程无指针", GOLD),
    ("④ 战斗 → 结算", "FHSRBattleResult → UHSRRewardSubsystem 生成经验 / 掉落 / 遗器", GREEN),
    ("⑤ 结算 → 存档", "各子系统 ExportSaveData → FHSRSaveData 快照 → 事务写入磁盘", GRAY),
]
dy = Inches(4.5)
for t, d, c in detail:
    add_rect(s, Inches(0.55), dy, Inches(12.23), Inches(0.52), fill=PANEL, line=PANEL2, radius=True)
    add_text(s, Inches(0.85), dy + Inches(0.05), Inches(2.4), Inches(0.4), t, size=13, color=c, bold=True)
    add_text(s, Inches(3.3), dy + Inches(0.05), Inches(9.3), Inches(0.42), d, size=12, color=GRAY)
    dy += Inches(0.6)
add_footer_note(s, "核心：权威结果永远是一份可回放、可存档、可测试的纯值数据，UI 只是它的投影。")

# ═══════════════════════════════════════════════════════════════════
# P14 各子系统关系图
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "SUBSYSTEMS", "各子系统关系 · 单一权威 + 纯值汇聚", 14)
add_text(s, Inches(0.55), Inches(1.4), Inches(5.5), Inches(0.35), "GameInstance 常驻 · 数据中心", size=15, color=CYAN, bold=True)
subsystems = [
    ("UHSRCharacterProfileSubsystem", "角色定义 / 成长 / 属性"),
    ("UHSRPartySubsystem", "队伍编成 / 切换"),
    ("UHSREquipmentSubsystem", "装备 · 遗器 · 套装 / 属性投影"),
    ("UHSRInventorySubsystem", "背包 / 物品 / 唯一实例"),
    ("UHSRRewardSubsystem", "掉落 / 奖励账本"),
    ("UHSRQuestSubsystem", "任务状态 / 目标事件"),
    ("UHSRMapSubsystem", "地图 / 传送 / 探索标志"),
    ("UHSRChallengeProgressionSubsystem", "挑战进度"),
]
sx = Inches(0.55); sw = Inches(3.95); sgap = Inches(0.1)
for i, (t, d) in enumerate(subsystems):
    col = i % 2; row = i // 2
    x = sx + (sw + sgap) * col
    y = Inches(1.8) + row * Inches(0.6)
    add_rect(s, x, y, sw, Inches(0.54), fill=PANEL, line=PANEL2, radius=True)
    add_text(s, x + Inches(0.15), y + Inches(0.05), Inches(1.9), Inches(0.45), t.replace('UHSR', ''), size=10.5, color=WHITE, bold=True)
    add_text(s, x + Inches(1.95), y + Inches(0.07), Inches(1.95), Inches(0.45), d, size=10.5, color=GRAY)
add_text(s, Inches(4.75), Inches(1.4), Inches(5.5), Inches(0.35), "存档是唯一汇聚点", size=15, color=GOLD, bold=True)
add_rect(s, Inches(4.75), Inches(1.8), Inches(8.05), Inches(2.3), fill=PANEL, line=GOLD, radius=True)
add_text(s, Inches(5.0), Inches(1.95), Inches(7.5), Inches(0.4), "UHSRSaveSubsystem", size=15, color=GOLD, bold=True)
add_text(s, Inches(5.0), Inches(2.35), Inches(7.5), Inches(0.5),
         "读所有子系统 → ExportSaveData → FHSRSaveData\n纯值 DTO 序列化 · 不保存任何 Actor / Widget 引用", size=12, color=GRAY)
add_text(s, Inches(5.0), Inches(3.15), Inches(7.5), Inches(0.4),
         "依赖关系：Reward → Inventory · Quest → Reward · Equipment 投影到角色 ASC", size=11, color=DIM)
add_text(s, Inches(4.75), Inches(4.35), Inches(5.5), Inches(0.35), "战斗世界 · 独立于探索", size=15, color=CYAN, bold=True)
add_rect(s, Inches(4.75), Inches(4.75), Inches(8.05), Inches(2.0), fill=PANEL, line=CYAN, radius=True)
add_text(s, Inches(5.0), Inches(4.9), Inches(7.5), Inches(0.4), "AHSRBattleGameMode → UHSRBattleCoordinator", size=13, color=CYAN, bold=True)
add_text(s, Inches(5.0), Inches(5.3), Inches(7.5), Inches(1.3),
         "· TurnManager 速度行动条（事件驱动，无 Tick）\n"
         "· Targeting 目标策略 · Status 状态系统\n"
         "· 与探索仅经 UHSRBattleTransitionSubsystem 握手\n"
         "  返回时携带 FHSRBattleResult 纯值结算", size=12, color=GRAY)
add_rect(s, Inches(0.55), Inches(6.35), Inches(12.23), Inches(0.75), fill=PANEL, line=GREEN, radius=True)
add_text(s, Inches(0.85), Inches(6.45), Inches(11.6), Inches(0.55),
         "两大世界（探索 / 战斗）各自完整，靠纯值请求与纯值结算对接；所有子系统最终在存档处归一。",
         size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_footer_note(s, "面试要点：能画清「谁拥有权威、数据在哪汇聚、边界传什么」比背类名更重要。")

# ═══════════════════════════════════════════════════════════════════
# P13 质量保障
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "QUALITY", "质量保障 · 每个阶段都可验证", 15)
# 流程
add_text(s, Inches(0.55), Inches(1.5), Inches(6), Inches(0.35), "TDD 流程（每个功能包）", size=15, color=CYAN, bold=True)
tdd = ["写失败测试\n(Red)", "实现\n(Green)", "构建 + 回归", "用户 PIE 验收"]
x = Inches(0.55); w = Inches(2.7); gap = Inches(0.15)
for i, lab in enumerate(tdd):
    add_pill(s, x, Inches(1.95), w, Inches(0.75), lab, [GREEN, BLUE, GOLD, CYAN][i], WHITE, 12)
    if i < 3:
        add_arrow(s, x + w + Inches(0.02), Inches(2.05), Inches(0.12), Inches(0.5))
    x += w + gap + Inches(0.1)
# 统计 tile
stats = [("30+", "自动化测试套件"), ("1000+", "构建持续通过"), ("18", "阶段逐一验收"), ("TDD", "先红后绿")]
sx = Inches(0.55); sw = Inches(3.0); sgap = Inches(0.1)
for i, (num, lab) in enumerate(stats):
    stat_tile(s, sx + (sw + sgap) * i, Inches(3.2), sw, Inches(1.5), num, lab, color=[CYAN, GOLD, GREEN, BLUE][i])
# 底部说明
add_card(s, Inches(0.55), Inches(5.1), Inches(12.2), Inches(1.6), "工程纪律",
         "· 每个功能包走真实 TDD：先有失败测试，再实现到通过\n· Development Editor 构建 + 聚焦自动化 + 相邻回归\n· 权威 / 表现 / 存档职责分离，表现失败可回退不影响权威结果\n· 用户 Editor/PIE 验收后才标记完成，证据分级记录",
         title_color=GOLD)
add_footer_note(s, "数字为示意口径，具体测试/构建计数按当前实际证据为准。")

# ═══════════════════════════════════════════════════════════════════
# P14 当前进度
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_header(s, "CURRENT STATUS", "当前进度", 16)
# 三栏
add_card(s, Inches(0.55), Inches(1.6), Inches(4.0), Inches(4.2), "已完成（可玩）",
         "探索 → 遭遇 → 回合战斗\n→ 结算 → 养成 → 存档\n\n完整闭环可玩\n全 UI 族（角色/队伍/背包/\n遗器/任务/地图/挑战/存档）\n18 阶段工程推进完成",
         title_color=GREEN)
add_card(s, Inches(4.8), Inches(1.6), Inches(4.0), Inches(4.2), "进行中",
         "正式 Demo 内容整合\n· 三张正式地图\n· 四名原创角色 / 12 技能\n· 两套遗器\n· 任务「界域回响」+ Boss\n\n资产校验门禁已 GREEN",
         title_color=CYAN)
add_card(s, Inches(9.05), Inches(1.6), Inches(3.75), Inches(4.2), "诚实边界",
         "· 部分美术为占位 / 灰盒\n· 角色立绘待作者化\n· 演示视频录制中\n\n边界如实说明，不虚报已完成。",
         title_color=GOLD)
add_text(s, Inches(0.55), Inches(6.1), Inches(12.2), Inches(0.5),
         "下一步：完成正式内容 → 录制分模块演示视频 → 成片用于作品集与讲解。", size=15, color=GRAY, bold=True)
add_footer_note(s, "")

# ═══════════════════════════════════════════════════════════════════
# P15 结尾
# ═══════════════════════════════════════════════════════════════════
s = new_slide()
add_text(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.2), "谢谢观看", size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(0.6),
         "每项功能配套一段录屏 / 截图讲解，正在逐步录制", size=18, color=CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.6),
         "探索  ·  战斗  ·  养成  ·  任务对话  ·  存档  ·  UI 全家桶", size=15, color=GRAY, align=PP_ALIGN.CENTER)
add_rect(s, Inches(3.4), Inches(5.2), Inches(6.5), Inches(0.06), fill=GOLD)
add_text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.5),
         "HSR · 从空工程到完整可玩 JRPG", size=16, color=GRAY, align=PP_ALIGN.CENTER)

prs.save(OUT_PATH)
print("saved:", OUT_PATH)
print("slides:", len(prs.slides))
